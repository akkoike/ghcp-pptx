"""
header.pptx (表紙) + PptxGenJS 出力 (コンテンツ) + footer.pptx (最終スライド) を結合し、
表紙のタイトルを変更して最終 .pptx を出力するスクリプト

SKILL.md に基づく copy_template_slide / copy_content_slide 方式を使用
"""
import os
import copy
from pptx import Presentation
from pptx.util import Pt

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ROOT_DIR = os.path.dirname(SCRIPT_DIR)
TEMPLATES_DIR = os.path.join(os.path.dirname(ROOT_DIR), 'templates')

HEADER_PPTX = os.path.join(TEMPLATES_DIR, 'header.pptx')
FOOTER_PPTX = os.path.join(TEMPLATES_DIR, 'footer.pptx')
CONTENT_PPTX = os.path.join(ROOT_DIR, 'docs', 'FDE_AI支援_content.pptx')
OUTPUT_PPTX = os.path.join(ROOT_DIR, 'docs', 'FDE_AI支援.pptx')

TITLE_TEXT = 'SIパートナーにおける FDE × AI支援の最新動向'

# ============================================================
# XML 名前空間
# ============================================================
NSMAP_P = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
NSMAP_R = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
NSMAP_A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
RT_SLIDE_LAYOUT = (
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout'
)
RT_SLIDE_MASTER = (
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster'
)
RT_NOTES_SLIDE = (
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide'
)
SKIP_RELTYPES = {RT_SLIDE_LAYOUT, RT_SLIDE_MASTER, RT_NOTES_SLIDE}


# ============================================================
# ヘルパー関数
# ============================================================
def _find_blank_layout(prs):
    """Blank レイアウトを探す"""
    for layout in prs.slide_layouts:
        if layout.name in ('Blank', '白紙'):
            return layout
    return prs.slide_layouts[6]  # フォールバック


def _collect_rels(part, rId_map, new_slide_part):
    """part のリレーションを new_slide_part にコピーし rId_map を更新する"""
    for rId, rel in part.rels.items():
        if rel.is_external:
            continue
        if rel.reltype in SKIP_RELTYPES:
            continue
        new_rId = new_slide_part.relate_to(rel.target_part, rel.reltype)
        rId_map[rId] = new_rId


def _get_ph_type(sp_elem):
    """shape 要素からプレースホルダータイプを取得する (None=非PH)"""
    nvPr = sp_elem.find(f'.//{NSMAP_P}nvPr')
    if nvPr is not None:
        ph = nvPr.find(f'{NSMAP_P}ph')
        if ph is not None:
            return ph.get('type', 'body')
    return None


def _resolve_placeholder_inheritance(slide_sp, layout_spTree):
    """スライドのプレースホルダーにレイアウトの spPr / bodyPr を解決する"""
    ph_type = _get_ph_type(slide_sp)
    if ph_type is None:
        return

    for layout_sp in layout_spTree:
        tag = layout_sp.tag.split('}')[-1] if '}' in layout_sp.tag else layout_sp.tag
        if tag in ('nvGrpSpPr', 'grpSpPr'):
            continue
        if _get_ph_type(layout_sp) == ph_type:
            # spPr が空なら、レイアウトの spPr で置換
            slide_spPr = slide_sp.find(f'{NSMAP_P}spPr')
            layout_spPr = layout_sp.find(f'{NSMAP_P}spPr')
            if slide_spPr is not None and layout_spPr is not None:
                if len(slide_spPr) == 0:
                    parent = slide_spPr.getparent()
                    idx = list(parent).index(slide_spPr)
                    parent.remove(slide_spPr)
                    parent.insert(idx, copy.deepcopy(layout_spPr))

            # bodyPr が空なら、レイアウトの bodyPr で置換
            slide_txBody = slide_sp.find(f'{NSMAP_P}txBody')
            layout_txBody = layout_sp.find(f'{NSMAP_P}txBody')
            if slide_txBody is not None and layout_txBody is not None:
                slide_bodyPr = slide_txBody.find(f'{NSMAP_A}bodyPr')
                layout_bodyPr = layout_txBody.find(f'{NSMAP_A}bodyPr')
                if slide_bodyPr is not None and layout_bodyPr is not None:
                    if len(slide_bodyPr) == 0 and len(slide_bodyPr.attrib) == 0:
                        parent = slide_bodyPr.getparent()
                        idx = list(parent).index(slide_bodyPr)
                        parent.remove(slide_bodyPr)
                        parent.insert(idx, copy.deepcopy(layout_bodyPr))
            break


# ============================================================
# copy_template_slide: テンプレートスライド用（レイアウト・フラットニング方式）
# ============================================================
def copy_template_slide(src_slide, dst_prs):
    """
    テンプレートスライド (header/footer) を、レイアウト由来のデザインを
    含めて正しくコピーする（レイアウト・フラットニング方式）。
    """
    src_layout = src_slide.slide_layout
    blank_layout = _find_blank_layout(dst_prs)
    new_slide = dst_prs.slides.add_slide(blank_layout)

    # placeholder を削除
    for ph in list(new_slide.placeholders):
        ph._element.getparent().remove(ph._element)

    # --- リレーションを両方からコピー ---
    rId_map = {}
    _collect_rels(src_layout.part, rId_map, new_slide.part)  # レイアウトから先
    _collect_rels(src_slide.part, rId_map, new_slide.part)   # スライドから後

    # --- cSld をディープコピー ---
    new_cSld = copy.deepcopy(
        src_slide._element.find(f'{NSMAP_P}cSld')
    )
    layout_cSld = src_layout._element.find(f'{NSMAP_P}cSld')

    # (A) レイアウトの背景をスライドに統合
    slide_bg = new_cSld.find(f'{NSMAP_P}bg')
    layout_bg = layout_cSld.find(f'{NSMAP_P}bg') if layout_cSld is not None else None
    if slide_bg is None and layout_bg is not None:
        bg_copy = copy.deepcopy(layout_bg)
        spTree = new_cSld.find(f'{NSMAP_P}spTree')
        spTree_idx = list(new_cSld).index(spTree) if spTree is not None else 0
        new_cSld.insert(spTree_idx, bg_copy)

    # (B) プレースホルダーの書式継承を解決
    slide_spTree = new_cSld.find(f'{NSMAP_P}spTree')
    layout_spTree = layout_cSld.find(f'{NSMAP_P}spTree') if layout_cSld is not None else None
    if slide_spTree is not None and layout_spTree is not None:
        slide_ph_types = set()
        for sp in list(slide_spTree):
            pt = _get_ph_type(sp)
            if pt is not None:
                slide_ph_types.add(pt)
                _resolve_placeholder_inheritance(sp, layout_spTree)

        # (C) レイアウトの非プレースホルダー shape をスライドに統合
        insert_pos = 2  # nvGrpSpPr, grpSpPr の後
        for child in layout_spTree:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag in ('nvGrpSpPr', 'grpSpPr'):
                continue
            if _get_ph_type(child) in slide_ph_types:
                continue
            slide_spTree.insert(insert_pos, copy.deepcopy(child))
            insert_pos += 1

    # --- 全 rId を書き換え ---
    for elem in new_cSld.iter():
        for attr_name, attr_val in list(elem.attrib.items()):
            if NSMAP_R in attr_name and attr_val in rId_map:
                elem.attrib[attr_name] = rId_map[attr_val]

    # --- cSld を差し替え ---
    old_cSld = new_slide._element.find(f'{NSMAP_P}cSld')
    new_slide._element.replace(old_cSld, new_cSld)

    return new_slide


# ============================================================
# copy_content_slide: PptxGenJS コンテンツスライド用
# ============================================================
def copy_content_slide(src_slide, dst_prs):
    """PptxGenJS 生成のコンテンツスライドをコピーする（rId マッピング方式）"""
    blank_layout = _find_blank_layout(dst_prs)
    new_slide = dst_prs.slides.add_slide(blank_layout)

    for ph in list(new_slide.placeholders):
        ph._element.getparent().remove(ph._element)

    rId_map = {}
    _collect_rels(src_slide.part, rId_map, new_slide.part)

    src_cSld = copy.deepcopy(
        src_slide._element.find(f'{NSMAP_P}cSld')
    )
    for elem in src_cSld.iter():
        for attr_name, attr_val in list(elem.attrib.items()):
            if NSMAP_R in attr_name and attr_val in rId_map:
                elem.attrib[attr_name] = rId_map[attr_val]

    old_cSld = new_slide._element.find(f'{NSMAP_P}cSld')
    new_slide._element.replace(old_cSld, src_cSld)
    return new_slide


# ============================================================
# メイン処理
# ============================================================
def main():
    # --- 1. header.pptx を開いてベースにする ---
    print(f'[python-pptx] header.pptx を読み込み中: {HEADER_PPTX}')
    prs = Presentation(HEADER_PPTX)

    # スライドサイズを 16:9 Wide に設定
    prs.slide_width = Pt(13.333 * 72)  # 13.333 inches
    prs.slide_height = Pt(7.5 * 72)    # 7.5 inches

    # 表紙のタイトルを変更
    cover_slide = prs.slides[0]
    for shape in cover_slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = ''.join([run.text for run in paragraph.runs])
                if 'Title' in full_text or 'title' in full_text or 'タイトル' in full_text:
                    for run in paragraph.runs:
                        run.text = ''
                    if paragraph.runs:
                        paragraph.runs[0].text = TITLE_TEXT
                        paragraph.runs[0].font.name = 'Meiryo UI'
                    break

    print(f'  表紙タイトルを「{TITLE_TEXT}」に変更しました')

    # --- 2. コンテンツスライド (PptxGenJS生成) を copy_content_slide で追加 ---
    print(f'[python-pptx] コンテンツスライドを読み込み中: {CONTENT_PPTX}')
    content_prs = Presentation(CONTENT_PPTX)
    for slide in content_prs.slides:
        copy_content_slide(slide, prs)
    print(f'  コンテンツスライド {len(content_prs.slides)} 枚を追加しました')

    # --- 3. footer.pptx を copy_template_slide（レイアウト・フラットニング方式）で追加 ---
    print(f'[python-pptx] footer.pptx を読み込み中: {FOOTER_PPTX}')
    footer_prs = Presentation(FOOTER_PPTX)
    for slide in footer_prs.slides:
        copy_template_slide(slide, prs)
    print(f'  footer スライド {len(footer_prs.slides)} 枚を追加しました')

    # --- 4. 保存 ---
    os.makedirs(os.path.dirname(OUTPUT_PPTX), exist_ok=True)
    prs.save(OUTPUT_PPTX)
    print(f'[python-pptx] 最終ファイルを出力しました: {OUTPUT_PPTX}')
    print(f'  スライド数: {len(prs.slides)}')


if __name__ == '__main__':
    main()
