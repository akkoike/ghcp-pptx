import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

# === Constants ===
FONT = "Meiryo UI"
BLUE = RGBColor(0x00, 0x33, 0x66)
ACCENT = RGBColor(0x00, 0x78, 0xD4)
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BLUE = RGBColor(0xDE, 0xEA, 0xF6)
LIGHT_GREEN = RGBColor(0xE2, 0xEF, 0xDA)
LIGHT_ORANGE = RGBColor(0xFC, 0xE4, 0xCC)
LIGHT_PURPLE = RGBColor(0xE8, 0xDA, 0xEF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
DARK_BLUE = RGBColor(0x1B, 0x4F, 0x72)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
BASE_DIR = os.path.dirname(SCRIPT_DIR)
TEMPLATE_DIR = os.path.join(os.path.dirname(BASE_DIR), "..", "templates")
OUTPUT_PATH = os.path.join(BASE_DIR, "docs", "富士通_攻めるべき領域分析.pptx")

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# === Helper Functions ===

def add_bg(slide, color=WHITE):
    bg = slide.background
    f = bg.fill
    f.solid()
    f.fore_color.rgb = color

def add_title_bar(slide, title_text, y=Inches(0.0), h=Inches(0.8)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), y, SLIDE_W, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = BLUE
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT
    p.alignment = PP_ALIGN.LEFT

def add_subtitle(slide, text, x, y, w, h, size=Pt(14), color=BLUE, bold=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = PP_ALIGN.LEFT

def add_textbox(slide, x, y, w, h, paragraphs, spacing=Pt(4)):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text = para.get('text', '')
        p.text = text
        p.font.size = para.get('size', Pt(11))
        p.font.bold = para.get('bold', False)
        p.font.color.rgb = para.get('color', TEXT_COLOR)
        p.font.name = FONT
        p.alignment = para.get('align', PP_ALIGN.LEFT)
        p.space_before = para.get('space_before', spacing)
        p.space_after = para.get('space_after', Pt(2))

def add_icon_box(slide, x, y, w, h, icon, title, bullets, bg_color=LIGHT_BLUE, icon_color=ACCENT):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg_color
    shp.line.fill.background()
    shp.shadow.inherit = False

    icon_shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), Inches(0.5), Inches(0.5))
    icon_shp.fill.solid()
    icon_shp.fill.fore_color.rgb = icon_color
    icon_shp.line.fill.background()
    tf_icon = icon_shp.text_frame
    tf_icon.word_wrap = False
    p_icon = tf_icon.paragraphs[0]
    p_icon.text = icon
    p_icon.font.size = Pt(16)
    p_icon.font.color.rgb = WHITE
    p_icon.font.name = FONT
    p_icon.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.75), w - Inches(0.3), h - Inches(0.85))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.font.name = FONT
    p.space_after = Pt(4)

    for b in bullets:
        p2 = tf.add_paragraph()
        p2.text = "• " + b
        p2.font.size = Pt(9)
        p2.font.color.rgb = TEXT_COLOR
        p2.font.name = FONT
        p2.space_before = Pt(2)
        p2.space_after = Pt(1)

def add_table(slide, x, y, w, h, headers, rows, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    table = table_shape.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT
            paragraph.alignment = PP_ALIGN.CENTER

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            bg = LIGHT_GRAY if r % 2 == 0 else WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = TEXT_COLOR
                paragraph.font.name = FONT
                paragraph.alignment = PP_ALIGN.LEFT

def add_page_number(slide, num, total):
    tb = slide.shapes.add_textbox(Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY
    p.font.name = FONT
    p.alignment = PP_ALIGN.RIGHT

def add_footer_line(slide):
    ln = slide.shapes.add_connector(1, Inches(0.3), Inches(7.0), Inches(13.0), Inches(7.0))
    ln.line.color.rgb = ACCENT
    ln.line.width = Pt(0.75)


# =======================================
# SLIDE 1: 表紙
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE)

tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "富士通から見た各社 生成AI 取り組みと"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = FONT
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "攻める価値のある領域分析"
p2.font.size = Pt(32)
p2.font.bold = True
p2.font.color.rgb = WHITE
p2.font.name = FONT
p2.alignment = PP_ALIGN.CENTER

tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.8))
tf2 = tb2.text_frame
tf2.word_wrap = True
p3 = tf2.paragraphs[0]
p3.text = "NEC・日立・NTT の最新動向と富士通の競争戦略"
p3.font.size = Pt(18)
p3.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
p3.font.name = FONT
p3.alignment = PP_ALIGN.CENTER

tb3 = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.5))
tf3 = tb3.text_frame
p4 = tf3.paragraphs[0]
p4.text = "2026年4月"
p4.font.size = Pt(14)
p4.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
p4.font.name = FONT
p4.alignment = PP_ALIGN.CENTER

ln = slide.shapes.add_connector(1, Inches(4.0), Inches(3.7), Inches(9.3), Inches(3.7))
ln.line.color.rgb = RGBColor(0x64, 0xB5, 0xF6)
ln.line.width = Pt(1.5)


# =======================================
# SLIDE 2: 目次
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "目次")

toc_items = [
    ("1", "富士通の生成AI戦略 ― 4つの柱"),
    ("2", "Fujitsu Kozuchi & Application Transform"),
    ("3", "Physical AI 1.0 & Enterprise AI Factory"),
    ("4", "独自LLM「Takane」の優位性"),
    ("5", "NEC の生成AI取り組み"),
    ("6", "日立製作所の生成AI取り組み"),
    ("7", "NTT の生成AI取り組み"),
    ("8", "4社比較サマリー"),
    ("9", "攻めるべき領域 ① 金融レガシー刷新"),
    ("10", "攻めるべき領域 ② 公共・自治体AI"),
    ("11", "攻めるべき領域 ③ 製造業フィジカルAI"),
    ("12", "攻めるべき領域 ④ AIエージェント自動化"),
    ("13", "攻めるべき領域 ⑤ AIインフラ省電力化"),
    ("14", "優先順位と推奨アクション"),
    ("15", "まとめ・提言"),
]

left_items = toc_items[:8]
right_items = toc_items[8:]

for col_items, x_start in [(left_items, Inches(0.8)), (right_items, Inches(6.8))]:
    for i, (num, title) in enumerate(col_items):
        y_pos = Inches(1.2) + Inches(i * 0.7)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_start, y_pos, Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        tf_c = circle.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.text = num
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = WHITE
        p_c.font.name = FONT
        p_c.alignment = PP_ALIGN.CENTER

        tb = slide.shapes.add_textbox(x_start + Inches(0.55), y_pos + Inches(0.05), Inches(5.0), Inches(0.35))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = FONT

add_footer_line(slide)


# =======================================
# SLIDE 3: 富士通の生成AI戦略 ― 4つの柱
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "富士通の生成AI戦略 ― 4つの柱")

pillars = [
    ("🔧", "Fujitsu Kozuchi", ["50以上のAI技術を提供する共創型AIプラットフォームです", "業務特化型AIエージェント（会議支援、規約解析等）を搭載しています", "顧客自ら技術検証が可能な環境を提供しています"], LIGHT_BLUE, ACCENT),
    ("📄", "Application Transform", ["COBOL等レガシーシステムの設計書自動生成サービスです", "作業量を約97%削減（1ヶ月→約1日）しました", "Knowledge Graph-Enhanced RAGで高精度を実現しています"], LIGHT_GREEN, GREEN),
    ("🤖", "Physical AI 1.0", ["NVIDIAとの協業によるロボット×AIエージェント連携です", "マルチAIエージェントフレームワークを構築しています", "機密性を担保しつつ業務自動化を推進しています"], LIGHT_ORANGE, ORANGE),
    ("🏭", "Enterprise AI Factory", ["オンプレミス環境でのセキュアなAI運用基盤です", "独自LLM「Takane」による日本語特化型AIです", "脆弱性診断・ガードレール等のAIセキュリティを装備しています"], LIGHT_PURPLE, RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (icon, title, bullets, bg_c, icon_c) in enumerate(pillars):
    x = Inches(0.3) + Inches(i * 3.2)
    add_icon_box(slide, x, Inches(1.2), Inches(3.0), Inches(4.5), icon, title, bullets, bg_c, icon_c)

add_textbox(slide, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.8), [
    {'text': '★ 富士通の差別化ポイント：「レガシー刷新力」×「Kozuchiプラットフォーム」×「セキュアなオンプレミスAI」', 'size': Pt(12), 'bold': True, 'color': RED}
])
add_footer_line(slide)


# =======================================
# SLIDE 4: Kozuchi & Application Transform
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Fujitsu Kozuchi & Application Transform")

add_subtitle(slide, "■ Fujitsu Kozuchi AIプラットフォーム", Inches(0.3), Inches(1.1), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.5), [
    {'text': '• 50以上のAI技術をWebアプリテンプレートで提供しています', 'size': Pt(10)},
    {'text': '• 会議支援、現場作業支援、規約解析などのAIエージェントを搭載しています', 'size': Pt(10)},
    {'text': '• テック主導の「共創型AIプラットフォーム」として位置づけられています', 'size': Pt(10)},
    {'text': '• 2026年に「Enterprise AI Factory」へ進化しました', 'size': Pt(10), 'bold': True, 'color': ACCENT},
])

add_subtitle(slide, "■ Application Transform（レガシー刷新）", Inches(0.3), Inches(4.0), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(4.5), Inches(5.8), Inches(2.5), [
    {'text': '• 2026年春からSaaS提供を開始しました', 'size': Pt(10)},
    {'text': '• 40年分のCOBOLコードを解析し設計書を自動生成します', 'size': Pt(10)},
    {'text': '• SMBC日興証券との共同検証を推進しています', 'size': Pt(10)},
])

ln_div = slide.shapes.add_connector(1, Inches(6.5), Inches(1.1), Inches(6.5), Inches(6.8))
ln_div.line.color.rgb = ACCENT
ln_div.line.width = Pt(1.5)

add_subtitle(slide, "■ Application Transform の成果指標", Inches(6.8), Inches(1.1), Inches(6.0), Inches(0.4))

metrics = [
    ("作業量削減", "約97%削減", "1ヶ月→約1日"),
    ("解析網羅性", "+95%向上", "仕様漏れを大幅削減"),
    ("設計書可読性", "+60%向上", "品質の均一化"),
    ("GitHub Copilot", "4,000人突破", "37.5万時間効率化"),
]

for i, (label, value, desc) in enumerate(metrics):
    y = Inches(1.7) + Inches(i * 1.3)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), y, Inches(6.0), Inches(1.1))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BLUE
    box.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(7.0), y + Inches(0.1), Inches(2.5), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.font.name = FONT

    tb2 = slide.shapes.add_textbox(Inches(9.5), y + Inches(0.05), Inches(3.0), Inches(0.45))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = value
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT
    p2.font.name = FONT

    tb3 = slide.shapes.add_textbox(Inches(7.0), y + Inches(0.55), Inches(5.5), Inches(0.35))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = desc
    p3.font.size = Pt(9)
    p3.font.color.rgb = GRAY
    p3.font.name = FONT

add_footer_line(slide)


# =======================================
# SLIDE 5: Physical AI & Enterprise AI Factory
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Physical AI 1.0 & Enterprise AI Factory")

add_subtitle(slide, "■ Fujitsu Kozuchi Physical AI 1.0", Inches(0.3), Inches(1.1), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(1.6), Inches(6.0), Inches(4.5), [
    {'text': '• NVIDIAとの協業により2025年末に発表しました', 'size': Pt(10)},
    {'text': '• AIエージェントとロボット等の物理デバイスが協調する「マルチAIエージェントフレームワーク」を構築しています', 'size': Pt(10)},
    {'text': '• 機密性を担保しつつ、業務自動化・協働ロボット連携を推進しています', 'size': Pt(10)},
    {'text': '• 今後はAIエージェントとロボットの自律進化・協調にも対応予定です', 'size': Pt(10)},
    {'text': '', 'size': Pt(6)},
    {'text': '【技術アーキテクチャ】', 'size': Pt(11), 'bold': True, 'color': BLUE},
    {'text': '• NVIDIA技術基盤上で動作するセキュアなAIエージェント群です', 'size': Pt(10)},
    {'text': '• 複数エージェントが物理デバイスと連携し複雑な業務を自動化します', 'size': Pt(10)},
    {'text': '• デジタルツイン技術との統合も視野に入れています', 'size': Pt(10)},
])

ln_div = slide.shapes.add_connector(1, Inches(6.5), Inches(1.1), Inches(6.5), Inches(6.8))
ln_div.line.color.rgb = ACCENT
ln_div.line.width = Pt(1.5)

add_subtitle(slide, "■ Enterprise AI Factory", Inches(6.8), Inches(1.1), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(4.5), [
    {'text': '• 2026年に発表されたオンプレミスAI運用基盤です', 'size': Pt(10)},
    {'text': '• 顧客の自社環境で機密データを安全に活用できます', 'size': Pt(10)},
    {'text': '• 業務特化AIエージェントの自動運用・管理を実現しています', 'size': Pt(10)},
    {'text': '', 'size': Pt(6)},
    {'text': '【セキュリティ機能】', 'size': Pt(11), 'bold': True, 'color': BLUE},
    {'text': '• 脆弱性診断機能を搭載しています', 'size': Pt(10)},
    {'text': '• ガードレール機能で出力品質を制御しています', 'size': Pt(10)},
    {'text': '• プロンプトインジェクション対策を実装しています', 'size': Pt(10)},
    {'text': '', 'size': Pt(6)},
    {'text': '【Takane LLM統合】', 'size': Pt(11), 'bold': True, 'color': BLUE},
    {'text': '• 継続的にアップデートできるファインチューニング機能を提供しています', 'size': Pt(10)},
    {'text': '• メモリ消費94%削減により効率的な運用が可能です', 'size': Pt(10)},
])

add_footer_line(slide)


# =======================================
# SLIDE 6: 独自LLM「Takane」
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "独自LLM「Takane」の優位性")

add_table(slide, Inches(0.5), Inches(1.2), Inches(5.5), Inches(3.5),
    ["項目", "内容"],
    [
        ["ベースモデル", "Cohere Command R+"],
        ["最適化領域", "日本語・業務用語に特化"],
        ["メモリ最適化", "94%削減（量子化技術）"],
        ["セキュリティ", "高セキュリティ設計"],
        ["カスタマイズ", "継続的ファインチューニング対応"],
        ["差別化", "日本企業の基幹システム・現場ニーズ対応"],
    ],
    [Inches(1.8), Inches(3.7)]
)

add_subtitle(slide, "■ 他社LLMとの比較", Inches(6.8), Inches(1.1), Inches(6.0), Inches(0.4))

add_table(slide, Inches(6.8), Inches(1.7), Inches(6.0), Inches(3.0),
    ["企業", "LLM名", "特徴"],
    [
        ["富士通", "Takane", "日本語業務特化、メモリ94%削減"],
        ["NEC", "cotomi", "日本語特化、暗黙知学習"],
        ["日立", "（他社LLM活用）", "OpenAI/Google連携"],
        ["NTT", "tsuzumi 2", "1GPU推論、省エネ設計"],
    ],
    [Inches(1.2), Inches(1.5), Inches(3.3)]
)

add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.5), Inches(1.5), [
    {'text': '★ 富士通 Takane の強み：Cohere Command R+ベースの高い基礎性能 ＋ 日本語・業務用語最適化 ＋ 94%メモリ削減による低コスト運用', 'size': Pt(12), 'bold': True, 'color': RED},
    {'text': '→ OpenAI / Google 等のクラウドLLMと比較し、「日本企業の基幹システム・レガシー刷新・現場ニーズ」に合わせたカスタム性が最大の差別化要素です', 'size': Pt(10), 'color': ACCENT},
])

add_footer_line(slide)


# =======================================
# SLIDE 7: NEC の取り組み
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "NEC の生成AI取り組み")

add_icon_box(slide, Inches(0.3), Inches(1.2), Inches(4.0), Inches(3.0),
    "🗣", "cotomi（LLM）",
    ["高精度な日本語処理能力を持つ国産LLMです", "RAG・ファインチューニングで業務特化が可能です", "セキュアなクローズド環境での運用に対応しています"],
    LIGHT_BLUE, ACCENT)

add_icon_box(slide, Inches(4.6), Inches(1.2), Inches(4.0), Inches(3.0),
    "🤖", "cotomi Act（エージェント）",
    ["Webブラウザ操作の暗黙知を学習・自動化します", "WebArena ベンチマークで80.4%を達成しました", "2026年度中にサービス提供を予定しています"],
    LIGHT_GREEN, GREEN)

add_icon_box(slide, Inches(8.9), Inches(1.2), Inches(4.0), Inches(3.0),
    "🔒", "AIガバナンス",
    ["ハルシネーション対策技術を開発しています", "Google Cloud・Cisco との協業を展開しています", "情報漏洩対策・権限管理を提供しています"],
    LIGHT_ORANGE, ORANGE)

add_subtitle(slide, "■ 主な導入事例", Inches(0.3), Inches(4.6), Inches(6.0), Inches(0.4))
add_table(slide, Inches(0.3), Inches(5.1), Inches(6.0), Inches(1.6),
    ["分野", "導入先", "活用内容"],
    [
        ["行政", "北九州市", "AI会計室による問い合わせ自動化"],
        ["行政", "相模原市・盛岡市", "議会答弁資料作成支援"],
        ["医療", "病院", "退院サマリー自動生成"],
    ],
    [Inches(1.0), Inches(2.0), Inches(3.0)]
)

add_subtitle(slide, "■ 注目ポイント", Inches(6.8), Inches(4.6), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(7.0), Inches(5.1), Inches(5.8), Inches(1.8), [
    {'text': '• PLM「Obbligato AI」で製造・設計分野にも展開しています', 'size': Pt(10)},
    {'text': '• 「クライアントゼロ」戦略で自社DXを先行実践しています', 'size': Pt(10)},
    {'text': '• GENIACプロジェクトで国産AI開発力強化に貢献しています', 'size': Pt(10)},
    {'text': '★ 公共・医療分野での導入実績が豊富です', 'size': Pt(10), 'bold': True, 'color': RED},
])

add_footer_line(slide)


# =======================================
# SLIDE 8: 日立の取り組み
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "日立製作所の生成AI取り組み")

add_icon_box(slide, Inches(0.3), Inches(1.2), Inches(4.0), Inches(3.0),
    "⚙", "Lumada 3.0",
    ["AI主導で社会インフラ変革を目指しています", "2025年度売上4兆円突破見込みです", "3年間で1.3兆円規模の投資を計画しています"],
    LIGHT_BLUE, ACCENT)

add_icon_box(slide, Inches(4.6), Inches(1.2), Inches(4.0), Inches(3.0),
    "🏗", "フィジカルAI + HMAX",
    ["NVIDIAと戦略提携し「HMAX」を共同開発しています", "鉄道・エネルギー・製造分野の現場自動化を推進中です", "2026年4月にフィジカルAI体験スタジオを開設しました"],
    LIGHT_GREEN, GREEN)

add_icon_box(slide, Inches(8.9), Inches(1.2), Inches(4.0), Inches(3.0),
    "👥", "全社AI活用",
    ["グループ28万人が生成AIを利用しています", "システム開発効率は3割向上しました", "AIアンバサダー制度で組織横断推進しています"],
    LIGHT_ORANGE, ORANGE)

add_subtitle(slide, "■ 技術的特徴", Inches(0.3), Inches(4.6), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(5.1), Inches(5.8), Inches(1.8), [
    {'text': '• IT × OT × AI の融合が最大の差別化要素です', 'size': Pt(10)},
    {'text': '• 熟練技術者の暗黙知をAIに取り込み技術継承を実現しています', 'size': Pt(10)},
    {'text': '• 保守AIの回答精度は9割まで向上しています', 'size': Pt(10)},
    {'text': '• 将来は「創造AI（第4世代AI）」への進化も展望しています', 'size': Pt(10)},
])

add_subtitle(slide, "■ 注目ポイント", Inches(6.8), Inches(4.6), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(7.0), Inches(5.1), Inches(5.8), Inches(1.8), [
    {'text': '• 百五銀行、塩野義製薬、三菱ケミカル等と業種別協創しています', 'size': Pt(10)},
    {'text': '• Google Cloud、Microsoftとの協業も拡大中です', 'size': Pt(10)},
    {'text': '★ 製造・鉄道・エネルギーの現場密着型AI活用が圧倒的に強いです', 'size': Pt(10), 'bold': True, 'color': RED},
])

add_footer_line(slide)


# =======================================
# SLIDE 9: NTT の取り組み
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "NTT の生成AI取り組み")

add_icon_box(slide, Inches(0.3), Inches(1.2), Inches(4.0), Inches(3.0),
    "💬", "tsuzumi 2（LLM）",
    ["2025年10月から提供を開始しています", "1GPUで推論可能な省エネ・低コスト設計です", "金融・医療・公共の専門知識を強化しています"],
    LIGHT_BLUE, ACCENT)

add_icon_box(slide, Inches(4.6), Inches(1.2), Inches(4.0), Inches(3.0),
    "💡", "IOWN光電融合",
    ["2026年にPEC-2を商用化予定です", "AIデータセンターの消費電力を最大1/8に削減します", "2032年にエネルギー効率100倍超を目指しています"],
    LIGHT_GREEN, GREEN)

add_icon_box(slide, Inches(8.9), Inches(1.2), Inches(4.0), Inches(3.0),
    "📊", "ビジネス成果",
    ["2025年Q1で相談案件1,800件に到達しました", "受注総額670億円に拡大しています", "2027年に年間5,000億円規模を見込んでいます"],
    LIGHT_ORANGE, ORANGE)

add_subtitle(slide, "■ AI For Quality Growth 戦略", Inches(0.3), Inches(4.6), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(5.1), Inches(5.8), Inches(1.8), [
    {'text': '• コンサル～実装～インフラ運用まで一気通貫で支援しています', 'size': Pt(10)},
    {'text': '• NVIDIA DGX B200を採用した省電力AI基盤を整備しています', 'size': Pt(10)},
    {'text': '• サイバーセキュリティ分野でのAI活用も強化中です', 'size': Pt(10)},
    {'text': '• AIコンステレーション（自律分散型AI連携）にも注力しています', 'size': Pt(10)},
])

add_subtitle(slide, "■ 注目ポイント", Inches(6.8), Inches(4.6), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(7.0), Inches(5.1), Inches(5.8), Inches(1.8), [
    {'text': '• 東京通信大学、富士フイルムBI等への導入実績があります', 'size': Pt(10)},
    {'text': '• Broadcom、Accton等との量産体制を確立しています', 'size': Pt(10)},
    {'text': '★ 通信インフラ×AIの光電融合技術は他社にない独自優位性です', 'size': Pt(10), 'bold': True, 'color': RED},
])

add_footer_line(slide)


# =======================================
# SLIDE 10: 4社比較サマリー
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "4社 生成AI取り組み比較サマリー")

add_table(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.0),
    ["領域", "富士通", "NEC", "日立", "NTT"],
    [
        ["独自LLM", "Takane\n(Cohere R+ベース)", "cotomi", "なし\n(他社LLM活用)", "tsuzumi 2"],
        ["AIエージェント", "Kozuchi\nAIエージェント", "cotomi Act\n(WebArena 80.4%)", "マルチAI\nエージェント", "AI Advisor"],
        ["フィジカルAI", "◎ Kozuchi\nPhysical AI 1.0", "△ 限定的", "◎ HMAX\n(NVIDIA連携)", "△ 限定的"],
        ["レガシー刷新", "◎ Application\nTransform (97%削減)", "△", "△", "△"],
        ["AIインフラ", "△", "△", "○ Lumada 3.0\n(4兆円規模)", "◎ IOWN光電融合\n(消費電力1/8)"],
        ["製造・OT", "○", "○ Obbligato AI", "◎ IT×OT融合", "△"],
        ["AIガバナンス", "○ ガードレール", "◎ Google/Cisco\n連携", "○ AIアンバサダー", "○"],
    ],
    [Inches(1.5), Inches(2.8), Inches(2.8), Inches(2.8), Inches(2.8)]
)

add_textbox(slide, Inches(0.3), Inches(6.4), Inches(12.7), Inches(0.5), [
    {'text': '◎=業界トップ級  ○=一定の実績  △=限定的・未注力', 'size': Pt(9), 'color': GRAY},
])

add_footer_line(slide)


# =======================================
# SLIDE 11: 攻めるべき領域① 金融レガシー刷新
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "攻めるべき領域 ① 金融業界向けレガシー刷新 × AI開発生産性革命")

star_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.1), Inches(2.0), Inches(0.5))
star_box.fill.solid()
star_box.fill.fore_color.rgb = RED
star_box.line.fill.background()
tf_s = star_box.text_frame
p_s = tf_s.paragraphs[0]
p_s.text = "優先度 ★★★"
p_s.font.size = Pt(14)
p_s.font.bold = True
p_s.font.color.rgb = WHITE
p_s.font.name = FONT
p_s.alignment = PP_ALIGN.CENTER

add_subtitle(slide, "■ 競合状況", Inches(0.3), Inches(1.9), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(2.4), Inches(6.0), Inches(1.5), [
    {'text': '• 各社とも金融業界への展開を進めていますが、「ソフトウェア開発の生産性100倍」を打ち出しているのは富士通だけです', 'size': Pt(10)},
    {'text': '• 金融機関はシステム更新の規模が大きく、レガシーCOBOL資産も膨大です', 'size': Pt(10)},
    {'text': '• Application Transform の最適ターゲットとなる市場です', 'size': Pt(10)},
])

add_subtitle(slide, "■ 攻めるべき理由", Inches(0.3), Inches(3.8), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(4.3), Inches(6.0), Inches(2.0), [
    {'text': '• 富士通のApplication Transformは COBOL設計書生成で97%の効率化を実現済みです', 'size': Pt(10), 'bold': True, 'color': ACCENT},
    {'text': '• Knowledge Graph-Enhanced RAGで高い解析精度を確保しています', 'size': Pt(10)},
    {'text': '• SMBC日興証券との共同検証で金融分野の実績があります', 'size': Pt(10)},
    {'text': '• 「2025年の崖」対策として最も実践的なソリューションです', 'size': Pt(10)},
])

ln_div = slide.shapes.add_connector(1, Inches(6.5), Inches(1.5), Inches(6.5), Inches(6.8))
ln_div.line.color.rgb = ACCENT
ln_div.line.width = Pt(1.5)

add_subtitle(slide, "■ 推奨アクション", Inches(6.8), Inches(1.9), Inches(6.0), Inches(0.4))
action_items = [
    "SMBC日興証券との検証成果を横展開し、メガバンク・地銀へアプローチします",
    "「レガシー刷新 × 開発生産性100倍」のメッセージを金融業界全体に訴求します",
    "Application Transform ＋ Kozuchiの統合ソリューションとして提案を強化します",
]
for i, item in enumerate(action_items):
    y = Inches(2.5) + Inches(i * 1.1)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), y, Inches(6.0), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GREEN
    box.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(7.0), y + Inches(0.1), Inches(5.6), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Action {i+1}: {item}"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_COLOR
    p.font.name = FONT
    p.font.bold = True if i == 0 else False

add_footer_line(slide)


# =======================================
# SLIDE 12: 攻めるべき領域② 公共・自治体
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "攻めるべき領域 ② 公共・自治体向けAI業務改革")

star_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.1), Inches(2.0), Inches(0.5))
star_box.fill.solid()
star_box.fill.fore_color.rgb = RED
star_box.line.fill.background()
tf_s = star_box.text_frame
p_s = tf_s.paragraphs[0]
p_s.text = "優先度 ★★★"
p_s.font.size = Pt(14)
p_s.font.bold = True
p_s.font.color.rgb = WHITE
p_s.font.name = FONT
p_s.alignment = PP_ALIGN.CENTER

add_subtitle(slide, "■ 競合状況", Inches(0.3), Inches(1.9), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(2.4), Inches(6.0), Inches(1.5), [
    {'text': '• NECが北九州市、相模原市、盛岡市など自治体導入で先行しています', 'size': Pt(10)},
    {'text': '• NTTも自治体向けクローズドAI（tsuzumi 2）で強みを持っています', 'size': Pt(10)},
    {'text': '• 自治体はレガシーシステムの刷新ニーズが極めて高い市場です', 'size': Pt(10)},
])

add_subtitle(slide, "■ 攻めるべき理由", Inches(0.3), Inches(3.8), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(4.3), Inches(6.0), Inches(2.0), [
    {'text': '• 自治体のレガシーシステムは COBOL資産が多く、Application Transformとの親和性が高いです', 'size': Pt(10), 'bold': True, 'color': ACCENT},
    {'text': '• Enterprise AI Factory によるオンプレミスAI運用は自治体の情報セキュリティ要件に適合します', 'size': Pt(10)},
    {'text': '• レガシー刷新 ＋ 業務AIの組み合わせはNEC・NTTにない独自提案が可能です', 'size': Pt(10)},
])

ln_div = slide.shapes.add_connector(1, Inches(6.5), Inches(1.5), Inches(6.5), Inches(6.8))
ln_div.line.color.rgb = ACCENT
ln_div.line.width = Pt(1.5)

add_subtitle(slide, "■ 推奨アクション", Inches(6.8), Inches(1.9), Inches(6.0), Inches(0.4))
action_items_2 = [
    "レガシー刷新＋業務AIの組み合わせで自治体DXパッケージを構成します",
    "自治体のCOBOL資産調査から入り、Application Transform導入を提案します",
    "Enterprise AI Factory のオンプレミス特性で情報セキュリティ要件をクリアします",
]
for i, item in enumerate(action_items_2):
    y = Inches(2.5) + Inches(i * 1.1)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), y, Inches(6.0), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GREEN
    box.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(7.0), y + Inches(0.1), Inches(5.6), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Action {i+1}: {item}"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_COLOR
    p.font.name = FONT

add_footer_line(slide)


# =======================================
# SLIDE 13: 攻めるべき領域③ 製造業フィジカルAI
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "攻めるべき領域 ③ 製造業向けフィジカルAIの拡大")

star_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.1), Inches(2.0), Inches(0.5))
star_box.fill.solid()
star_box.fill.fore_color.rgb = ORANGE
star_box.line.fill.background()
tf_s = star_box.text_frame
p_s = tf_s.paragraphs[0]
p_s.text = "優先度 ★★☆"
p_s.font.size = Pt(14)
p_s.font.bold = True
p_s.font.color.rgb = WHITE
p_s.font.name = FONT
p_s.alignment = PP_ALIGN.CENTER

add_subtitle(slide, "■ 競合状況", Inches(0.3), Inches(1.9), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(2.4), Inches(6.0), Inches(1.5), [
    {'text': '• 日立がLumada 3.0 + HMAX で製造・鉄道・エネルギー分野を席巻しています', 'size': Pt(10)},
    {'text': '• 3年で1.3兆円規模の投資を計画、フィジカルAI体験スタジオも開設しました', 'size': Pt(10)},
    {'text': '• 日立のIT × OT融合力と現場知見は大きなアドバンテージです', 'size': Pt(10)},
])

add_subtitle(slide, "■ 攻めるべき理由", Inches(0.3), Inches(3.8), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(4.3), Inches(6.0), Inches(2.0), [
    {'text': '• Kozuchi Physical AI 1.0 は発表済みだが、実績面では日立に差があります', 'size': Pt(10)},
    {'text': '• 製造業DXの巨大市場は成長余地が大きく、複数プレイヤーが共存可能です', 'size': Pt(10), 'bold': True, 'color': ACCENT},
    {'text': '• レガシー刷新（Application Transform）と組み合わせた End-to-End提案で差別化できます', 'size': Pt(10)},
])

ln_div = slide.shapes.add_connector(1, Inches(6.5), Inches(1.5), Inches(6.5), Inches(6.8))
ln_div.line.color.rgb = ACCENT
ln_div.line.width = Pt(1.5)

add_subtitle(slide, "■ 推奨アクション", Inches(6.8), Inches(1.9), Inches(6.0), Inches(0.4))
action_items_3 = [
    "NVIDIA連携をさらに深化し、自動車・半導体等の日立が手薄な製造分野に集中投資します",
    "Application Transform ＋ Physical AI の End-to-End 提案を構成します",
    "先行する日立の「鉄道・エネルギー」を避け、ニッチ製造分野から実績を積みます",
]
for i, item in enumerate(action_items_3):
    y = Inches(2.5) + Inches(i * 1.1)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), y, Inches(6.0), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_ORANGE
    box.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(7.0), y + Inches(0.1), Inches(5.6), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Action {i+1}: {item}"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_COLOR
    p.font.name = FONT

add_footer_line(slide)


# =======================================
# SLIDE 14: 攻めるべき領域④ AIエージェント
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "攻めるべき領域 ④ AIエージェントによるWebブラウザ業務自動化")

star_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.1), Inches(2.0), Inches(0.5))
star_box.fill.solid()
star_box.fill.fore_color.rgb = ORANGE
star_box.line.fill.background()
tf_s = star_box.text_frame
p_s = tf_s.paragraphs[0]
p_s.text = "優先度 ★★☆"
p_s.font.size = Pt(14)
p_s.font.bold = True
p_s.font.color.rgb = WHITE
p_s.font.name = FONT
p_s.alignment = PP_ALIGN.CENTER

add_subtitle(slide, "■ 競合状況", Inches(0.3), Inches(1.9), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(2.4), Inches(6.0), Inches(1.5), [
    {'text': '• NECのcotomi ActがWebArena 80.4%（人間超え）で先行しています', 'size': Pt(10)},
    {'text': '• 暗黙知のデータ化とWebブラウザ自動操作を組み合わせた独自技術です', 'size': Pt(10)},
    {'text': '• 2026年度中にWebブラウザ拡張としてサービス提供予定です', 'size': Pt(10)},
])

add_subtitle(slide, "■ 攻めるべき理由", Inches(0.3), Inches(3.8), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(4.3), Inches(6.0), Inches(2.0), [
    {'text': '• KozuchiプラットフォームにWeb業務自動化機能を統合すれば、エンタープライズ全体の自動化を包括提供可能です', 'size': Pt(10), 'bold': True, 'color': ACCENT},
    {'text': '• Enterprise AI Factory との統合でオンプレミス環境でのセキュアな業務自動化が実現できます', 'size': Pt(10)},
    {'text': '• AIエージェント市場は急成長中で、先行投資の価値が高いです', 'size': Pt(10)},
])

ln_div = slide.shapes.add_connector(1, Inches(6.5), Inches(1.5), Inches(6.5), Inches(6.8))
ln_div.line.color.rgb = ACCENT
ln_div.line.width = Pt(1.5)

add_subtitle(slide, "■ 推奨アクション", Inches(6.8), Inches(1.9), Inches(6.0), Inches(0.4))
action_items_4 = [
    "KozuchiエージェントにWebブラウザ業務自動化機能を追加開発します",
    "Enterprise AI Factory との統合でオンプレミスでのセキュアな自動化を差別化します",
    "MCP（Model Context Protocol）対応で外部ITサービスとの連携を強化します",
]
for i, item in enumerate(action_items_4):
    y = Inches(2.5) + Inches(i * 1.1)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), y, Inches(6.0), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BLUE
    box.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(7.0), y + Inches(0.1), Inches(5.6), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Action {i+1}: {item}"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_COLOR
    p.font.name = FONT

add_footer_line(slide)


# =======================================
# SLIDE 15: 攻めるべき領域⑤ AIインフラ
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "攻めるべき領域 ⑤ AIインフラ省電力化")

star_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.1), Inches(2.0), Inches(0.5))
star_box.fill.solid()
star_box.fill.fore_color.rgb = GRAY
star_box.line.fill.background()
tf_s = star_box.text_frame
p_s = tf_s.paragraphs[0]
p_s.text = "優先度 ★☆☆"
p_s.font.size = Pt(14)
p_s.font.bold = True
p_s.font.color.rgb = WHITE
p_s.font.name = FONT
p_s.alignment = PP_ALIGN.CENTER

add_subtitle(slide, "■ 競合状況", Inches(0.3), Inches(1.9), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(2.4), Inches(6.0), Inches(1.5), [
    {'text': '• NTTがIOWN光電融合技術で圧倒的な技術優位性を持っています', 'size': Pt(10)},
    {'text': '• 2026年にPEC-2を商用化し、消費電力を最大1/8に削減します', 'size': Pt(10)},
    {'text': '• Broadcom、Accton等との量産体制も確立済みです', 'size': Pt(10)},
])

add_subtitle(slide, "■ 攻めるべき理由", Inches(0.3), Inches(3.8), Inches(6.0), Inches(0.4))
add_textbox(slide, Inches(0.5), Inches(4.3), Inches(6.0), Inches(2.0), [
    {'text': '• AI利用コスト・消費電力が企業のAI導入障壁であり、解決の意義は大きいです', 'size': Pt(10)},
    {'text': '• ただしNTTの技術優位性は極めて高く、正面から競合するのは非効率です', 'size': Pt(10), 'color': RED},
    {'text': '• 提携戦略がより現実的で、自社の強み（Takaneメモリ94%削減等）で補完可能です', 'size': Pt(10)},
])

ln_div = slide.shapes.add_connector(1, Inches(6.5), Inches(1.5), Inches(6.5), Inches(6.8))
ln_div.line.color.rgb = ACCENT
ln_div.line.width = Pt(1.5)

add_subtitle(slide, "■ 推奨アクション", Inches(6.8), Inches(1.9), Inches(6.0), Inches(0.4))
action_items_5 = [
    "NTTのIOWN技術との戦略提携を検討し、Kozuchiの基盤として取り込みます",
    "Takaneのメモリ94%削減技術をさらに発展させ、推論コスト削減で差別化します",
    "「AIモデル効率化」の領域でソフトウェア面からインフラ最適化に貢献します",
]
for i, item in enumerate(action_items_5):
    y = Inches(2.5) + Inches(i * 1.1)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), y, Inches(6.0), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(7.0), y + Inches(0.1), Inches(5.6), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Action {i+1}: {item}"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_COLOR
    p.font.name = FONT

add_footer_line(slide)


# =======================================
# SLIDE 16: 優先順位まとめ
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "攻めるべき領域の優先順位と推奨アクション")

add_table(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(3.5),
    ["優先度", "領域", "主な競合", "富士通の強み", "推奨アクション"],
    [
        ["★★★", "金融レガシー刷新\n＋AI開発生産性革命", "各社（富士通\n独自の強み）", "Application Transform\n(97%削減)", "SMBC日興検証成果の横展開\nメガバンク・地銀へアプローチ"],
        ["★★★", "公共・自治体向け\nAI業務改革", "NEC、NTT", "レガシー刷新＋\n業務AI統合提案", "自治体DXパッケージの構成\nCOBOL資産調査から入る"],
        ["★★☆", "製造業向け\nフィジカルAI拡大", "日立", "Physical AI 1.0\n＋NVIDIA連携", "自動車・半導体等に集中投資\nEnd-to-End提案で差別化"],
        ["★★☆", "AIエージェント\nWeb自動化", "NEC", "Kozuchiプラットフォーム\n統合力", "Web業務自動化機能の追加\nMCP対応強化"],
        ["★☆☆", "AIインフラ\n省電力化", "NTT", "Takaneメモリ\n94%削減", "IOWN提携の検討\nソフトウェア面のインフラ最適化"],
    ],
    [Inches(1.0), Inches(2.4), Inches(2.0), Inches(2.8), Inches(4.5)]
)

add_textbox(slide, Inches(0.3), Inches(5.0), Inches(12.7), Inches(1.8), [
    {'text': '【戦略の方向性】', 'size': Pt(13), 'bold': True, 'color': BLUE},
    {'text': '• 最優先: 金融・公共分野での「レガシー刷新 × AI業務改革」を富士通独自の武器として攻めます', 'size': Pt(11), 'bold': True, 'color': RED},
    {'text': '• 並行推進: 日立のフィジカルAI・NECのAIエージェントに対抗する機能強化を進めます', 'size': Pt(11)},
    {'text': '• 補完戦略: NTTのAIインフラ技術とは提携による補完関係を構築します', 'size': Pt(11)},
])

add_footer_line(slide)


# =======================================
# SLIDE 17: まとめ・提言
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "まとめ・提言")

add_subtitle(slide, "■ 富士通の最大の差別化ポイント", Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.4))

diff_items = [
    ("レガシーシステム刷新", "Application Transform による COBOL 設計書自動生成（97%効率化）は各社に対して独自の優位性があります"),
    ("Kozuchi 統合プラットフォーム", "50以上のAI技術 ＋ エージェント ＋ Physical AI を一気通貫で提供できる包括力があります"),
    ("セキュアなオンプレミスAI", "Enterprise AI Factory による自社環境での安全なAI運用は、機密性の高い業界に最適です"),
]

for i, (title, desc) in enumerate(diff_items):
    y = Inches(1.8) + Inches(i * 1.2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), y, Inches(12.7), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BLUE
    box.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), y + Inches(0.1), Inches(12.3), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"✓ {title}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.font.name = FONT
    tb2 = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.5), Inches(12.0), Inches(0.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_COLOR
    p2.font.name = FONT

add_subtitle(slide, "■ 推奨する戦略的方向性", Inches(0.3), Inches(5.5), Inches(12.7), Inches(0.4))

conclusion_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.8))
conclusion_box.fill.solid()
conclusion_box.fill.fore_color.rgb = RGBColor(0xFC, 0xE4, 0xCC)
conclusion_box.line.fill.background()
tb_c = slide.shapes.add_textbox(Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.7))
tf_c = tb_c.text_frame
tf_c.word_wrap = True
p_c = tf_c.paragraphs[0]
p_c.text = "金融・公共分野でのレガシー刷新 × AI業務改革を最優先で攻め、日立のフィジカルAI・NECのAIエージェントには機能強化で対抗し、NTTのAIインフラ技術とは提携で補完する三層戦略を推奨します"
p_c.font.size = Pt(12)
p_c.font.bold = True
p_c.font.color.rgb = RED
p_c.font.name = FONT

add_footer_line(slide)


# =======================================
# Add page numbers
# =======================================
TOTAL_SLIDES = len(prs.slides)
for i, slide in enumerate(prs.slides):
    if i == 0:
        continue
    add_page_number(slide, i + 1, TOTAL_SLIDES)


# =======================================
# Insert header template if available
# =======================================
header_path = os.path.normpath(os.path.join(TEMPLATE_DIR, "header.pptx"))
footer_path = os.path.normpath(os.path.join(TEMPLATE_DIR, "footer.pptx"))

if os.path.exists(header_path):
    try:
        header_prs = Presentation(header_path)
        if len(header_prs.slides) > 0:
            header_slide = header_prs.slides[0]
            cover_slide = prs.slides[0]
            for shape in header_slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if "Title" in para.text or "title" in para.text:
                            para.text = "富士通から見た各社 生成AI 取り組みと攻める価値のある領域分析"
                            for run in para.runs:
                                run.font.name = FONT
            print(f"Header template found: {header_path}")
    except Exception as e:
        print(f"Header template error (continuing without): {e}")

if os.path.exists(footer_path):
    try:
        footer_prs = Presentation(footer_path)
        if len(footer_prs.slides) > 0:
            print(f"Footer template found: {footer_path}")
    except Exception as e:
        print(f"Footer template error (continuing without): {e}")


# =======================================
# Save
# =======================================
os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
prs.save(OUTPUT_PATH)
print(f"\n✅ PowerPoint saved: {OUTPUT_PATH}")
print(f"   Total slides: {TOTAL_SLIDES}")
