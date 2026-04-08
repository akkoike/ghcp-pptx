"""
header.pptx (表紙) + PptxGenJS 出力 (コンテンツ) + footer.pptx (最終スライド) を結合し、
表紙のタイトルを「Azure Virtual Desktop 詳細仕様」に変更して最終 .pptx を出力するスクリプト
"""
import os
import copy
from pptx import Presentation
from pptx.util import Pt

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ROOT_DIR = os.path.dirname(SCRIPT_DIR)
TEMPLATES_DIR = os.path.join(os.path.dirname(ROOT_DIR), 'templates')

HEADER_PPTX = os.path.join(TEMPLATES_DIR, 'header.pptx')
FOOTER_PPTX = os.path.join(TEMPLATES_DIR, 'footer.pptx')
CONTENT_PPTX = os.path.join(ROOT_DIR, 'docs', 'AzureVirtualDesktop_content.pptx')
OUTPUT_PPTX = os.path.join(ROOT_DIR, 'docs', 'AzureVirtualDesktop.pptx')

TITLE_TEXT = 'Azure Virtual Desktop 詳細仕様'


def copy_slide(src_prs, src_slide, dst_prs):
    """src_slide を dst_prs の末尾にコピーする"""
    slide_layout = dst_prs.slide_layouts[6]  # Blank layout
    new_slide = dst_prs.slides.add_slide(slide_layout)

    # 既存の placeholders を削除
    for ph in list(new_slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # ソーススライドの要素をコピー
    for shape in src_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    # 背景をコピー
    if src_slide.background and src_slide.background._element is not None:
        bg = copy.deepcopy(src_slide.background._element)
        new_slide_element = new_slide._element
        existing_bg = new_slide_element.find(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}bg'
        )
        if existing_bg is not None:
            new_slide_element.remove(existing_bg)
        cSld = new_slide_element.find(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
        )
        if cSld is not None:
            cSld.insert(0, bg)

    return new_slide


def main():
    # --- 1. header.pptx を開いてベースにする ---
    prs = Presentation(HEADER_PPTX)

    # スライドサイズを 16:9 Wide に設定
    prs.slide_width = Pt(13.333 * 72)  # 13.333 inches
    prs.slide_height = Pt(7.5 * 72)    # 7.5 inches

    # 表紙のタイトルを変更
    cover_slide = prs.slides[0]
    for shape in cover_slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = ''.join([run.text for run in paragraph.runs])
                if 'Title' in full_text or 'title' in full_text or 'タイトル' in full_text:
                    for run in paragraph.runs:
                        run.text = ''
                    if paragraph.runs:
                        paragraph.runs[0].text = TITLE_TEXT
                        paragraph.runs[0].font.name = 'Meiryo UI'
                    break

    # --- 2. コンテンツスライドを追加 ---
    content_prs = Presentation(CONTENT_PPTX)
    for slide in content_prs.slides:
        copy_slide(content_prs, slide, prs)

    # --- 3. footer.pptx を追加 ---
    footer_prs = Presentation(FOOTER_PPTX)
    for slide in footer_prs.slides:
        copy_slide(footer_prs, slide, prs)

    # --- 4. 保存 ---
    os.makedirs(os.path.dirname(OUTPUT_PPTX), exist_ok=True)
    prs.save(OUTPUT_PPTX)
    print(f'[python-pptx] 最終ファイルを出力しました: {OUTPUT_PPTX}')
    print(f'  スライド数: {len(prs.slides)}')


if __name__ == '__main__':
    main()
